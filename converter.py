from pdf2docx import Converter
import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
import os


# Function to convert PDF to DOCX
def pdf_to_docx(pdf_path):
    docx_path = pdf_path.replace('.pdf', '.docx')
    cv = Converter(pdf_path)
    cv.convert(docx_path, start=0, end=None)
    cv.close()
    return docx_path


# Function to convert DOCX to PDF
def docx_to_pdf(docx_path):
    from fpdf import FPDF
    doc = Document(docx_path)
    pdf_path = docx_path.replace('.docx', '.pdf')
    
    pdf = FPDF()
    pdf.add_page()
    
    # Iterate through the paragraphs in the DOCX
    for para in doc.paragraphs:
        pdf.set_font("Arial", size=12)
        pdf.multi_cell(0, 10, para.text)

    pdf.output(pdf_path)
    return pdf_path


# Function to convert PDF to TXT
def pdf_to_txt(pdf_path):
    # Open the PDF file
    doc = fitz.open(pdf_path)  
    text = ""
    
    # Iterate through all the pages in the PDF
    for page_num in range(doc.page_count):
        page = doc.load_page(page_num)  # Load the page
        text += page.get_text()  # Extract text from the page
    
    # Save the extracted text into a .txt file
    txt_path = pdf_path.replace('.pdf', '.txt')
    with open(txt_path, 'w', encoding='utf-8') as f:
        f.write(text)
    
    return txt_path


# Function to convert DOCX to TXT
def docx_to_txt(docx_path):
    doc = Document(docx_path)
    txt_path = docx_path.replace('.docx', '.txt')
    
    # Write each paragraph from the DOCX to a .txt file
    with open(txt_path, 'w', encoding='utf-8') as f:
        for para in doc.paragraphs:
            f.write(para.text + '\n')
    
    return txt_path


# Function to convert PPTX to PDF
def pptx_to_pdf(pptx_path):
    pptx = Presentation(pptx_path)
    pdf_path = pptx_path.replace('.pptx', '.pdf')
    
    from fpdf import FPDF
    pdf = FPDF()
    pdf.add_page()
    
    # Iterate through each slide and extract text
    for slide in pptx.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                pdf.set_font("Arial", size=12)
                pdf.multi_cell(0, 10, shape.text)
    
    pdf.output(pdf_path)
    return pdf_path
